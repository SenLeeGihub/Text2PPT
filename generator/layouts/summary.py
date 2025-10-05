from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from .utils import ensure_bg
from ..ppt_builder import load_theme, add_title

def render(prs, routed_content):
    theme = load_theme()
    for s in routed_content.slides:
        add_title(prs, theme, s["title"])
        slide = prs.slides[-1]
        ensure_bg(slide, theme)

        left = theme.margins["left"]
        top = theme.margins["top"] + Inches(1.2)
        width = Inches(10) - theme.margins["left"] - theme.margins["right"]
        height = Inches(5.5)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True

        for i, b in enumerate(s["bullets"]):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"• {b}"
            p.font.size = Pt(theme.sizes["body_pt"])
            p.font.name = theme.fonts["body"]
            p.font.color.rgb = theme.colors["text"]
            p.space_after = Pt(8)   # 段落间距
            p.alignment = PP_ALIGN.LEFT
