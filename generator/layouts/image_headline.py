from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from ..ppt_builder import load_theme, add_title
from .utils import ensure_bg
from pathlib import Path

def render(prs, routed_content):
    theme = load_theme()
    for s in routed_content.slides:
        add_title(prs, theme, s["title"])
        slide = prs.slides[-1]
        ensure_bg(slide, theme)

        # 图片
        img = s.get("hero_image")
        if img and Path(img).exists():
            slide.shapes.add_picture(img, Inches(0.8), Inches(1.8), width=Inches(8.8))  # add_picture API :contentReference[oaicite:10]{index=10}

        # 说明
        if s.get("caption"):
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(8.8), Inches(1.0))
            p = tb.text_frame.paragraphs[0]; p.text = s["caption"]
            p.font.size = Pt(theme.sizes["body_pt"]); p.font.name = theme.fonts["body"]
