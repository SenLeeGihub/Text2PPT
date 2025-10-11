from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MsoShape, MSO_CONNECTOR_TYPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

from .utils import ensure_bg
from ..ppt_builder import load_theme, add_title


def _to_emu(value):
    if isinstance(value, int):
        return value
    return int(round(float(value)))


def _apply_cell_border(cell, color_hex="000000"):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tcPr.find(qn(tag))
        if existing is not None:
            tcPr.remove(existing)
        ln = OxmlElement(tag)
        ln.set("w", str(int(Pt(1).emu)))

        solid_fill = OxmlElement("a:solidFill")
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", color_hex)
        solid_fill.append(srgb)
        ln.append(solid_fill)

        dash = OxmlElement("a:prstDash")
        dash.set("val", "solid")
        ln.append(dash)

        tcPr.append(ln)


def render(prs, routed_content):
    theme = load_theme()
    margins = theme.margins
    left_margin = margins.get("left", Inches(0.6))
    right_margin = margins.get("right", Inches(0.6))
    top_margin = margins.get("top", Inches(0.6))
    bottom_margin = margins.get("bottom", Inches(0.6))

    drawing_width = prs.slide_width - left_margin - right_margin
    spine_thickness = Inches(0.05)
    max_events = 8
    box_max_width = Inches(3.0)
    box_height = Inches(1.8)
    connector_length = Inches(0.3)
    node_size = Inches(0.22)
    panel_gap = Inches(0.3)
    min_left_panel = Inches(2.5)
    min_right_panel = Inches(2.0)

    for slide_spec in routed_content.slides:
        seen_dates: set[str] = set()
        title_text = (
            slide_spec.get("title")
            or slide_spec.get("header_title")
            or slide_spec.get("heading")
            or "时间线"
        )
        slide = add_title(prs, theme, title_text)
        ensure_bg(slide, theme)

        events = (slide_spec.get("events") or [])[:max_events]
        if not events:
            continue

        left = left_margin
        right = left + drawing_width
        spine_y = top_margin + int(prs.slide_height * 0.28)

        spine = slide.shapes.add_shape(
            MsoShape.RECTANGLE,
            left,
            spine_y - spine_thickness / 2,
            drawing_width,
            spine_thickness,
        )
        spine.fill.solid()
        spine.fill.fore_color.rgb = theme.colors["sub"]
        spine.line.fill.background()

        tip = slide.shapes.add_shape(
            MsoShape.ISOSCELES_TRIANGLE,
            right - Inches(0.35),
            spine_y - Inches(0.18),
            Inches(0.35),
            Inches(0.36),
        )
        tip.fill.solid()
        tip.fill.fore_color.rgb = theme.colors["sub"]
        tip.line.fill.background()
        tip.rotation = 90

        event_count = len(events)
        span = drawing_width / (event_count - 1) if event_count > 1 else 0

        for idx, ev in enumerate(events):
            cx = left + (span * idx if event_count > 1 else drawing_width / 2)
            cx = int(cx)

            node = slide.shapes.add_shape(
                MsoShape.DIAMOND,
                cx - node_size / 2,
                spine_y - node_size / 2,
                node_size,
                node_size,
            )
            node.fill.solid()
            node.fill.fore_color.rgb = theme.colors["accent"]
            node.line.fill.background()

            usable_span = float(span) if event_count > 1 else float(drawing_width)
            box_width = int(min(box_max_width, usable_span * 0.85))
            box_width = max(box_width, int(Inches(1.8)))
            direction = -1 if idx % 2 == 0 else 1

            preferred_center = cx + int(direction * (box_width / 2))
            min_center = left + box_width // 2
            max_center = right - box_width // 2
            box_center = max(min_center, min(max_center, preferred_center))
            box_left = box_center - box_width // 2

            end_x = box_center
            end_y = spine_y + direction * connector_length
            connector = slide.shapes.add_connector(
                MSO_CONNECTOR_TYPE.STRAIGHT,
                _to_emu(cx),
                _to_emu(spine_y),
                _to_emu(end_x),
                _to_emu(end_y),
            )
            connector.line.color.rgb = theme.colors["sub"]
            connector.line.width = Pt(1.5)

            date_text = (ev.get("date") or "").strip()
            if date_text:
                display_date = date_text
                suffix = 2
                while display_date in seen_dates:
                    display_date = f"{date_text}-事件{suffix}"
                    suffix += 1
                seen_dates.add(display_date)

                date_width = Inches(1.5)
                date_box = slide.shapes.add_textbox(
                    cx - date_width / 2,
                    spine_y + Inches(0.08),
                    date_width,
                    Inches(0.35),
                )
                dtf = date_box.text_frame
                dtf.text = display_date
                dp = dtf.paragraphs[0]
                dp.font.size = Pt(12)
                dp.font.name = theme.fonts["body"]
                dp.font.color.rgb = theme.colors["sub"]
                dp.alignment = PP_ALIGN.CENTER

            box_top = end_y - box_height + Inches(0.05) if direction < 0 else end_y + Inches(0.05)
            text_box = slide.shapes.add_textbox(box_left, box_top, box_width, box_height)
            tf = text_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.BOTTOM if direction < 0 else MSO_ANCHOR.TOP

            headline = (ev.get("headline") or "").strip()
            detail = (ev.get("detail") or "").strip()

            p_headline = tf.paragraphs[0]
            p_headline.text = headline
            p_headline.font.size = Pt(12)
            p_headline.font.bold = True
            p_headline.font.name = theme.fonts["body"]
            p_headline.font.color.rgb = theme.colors["text"]
            p_headline.alignment = PP_ALIGN.CENTER

            if detail:
                p_detail = tf.add_paragraph()
                p_detail.text = detail
                p_detail.font.size = Pt(12)
                p_detail.font.name = theme.fonts["body"]
                p_detail.font.color.rgb = theme.colors["sub"]
                p_detail.alignment = PP_ALIGN.CENTER

        analysis = slide_spec.get("analysis") or {}
        table_headers = [h.strip() for h in (analysis.get("table_headers") or []) if h and h.strip()][:4]
        raw_rows = analysis.get("table_rows") or []
        table_rows = []
        for row in raw_rows:
            if not isinstance(row, list):
                continue
            cells = [str(c).strip() for c in row]
            if any(cells):
                table_rows.append(cells)

        if table_headers:
            col_count = len(table_headers)
        elif table_rows:
            col_count = len(table_rows[0])
            table_headers = [f"列{i + 1}" for i in range(col_count)]
        else:
            col_count = 2
            table_headers = ["时间", "要点"]

        normalized_rows = []
        for row in table_rows:
            padded = row + [""] * max(0, col_count - len(row))
            normalized_rows.append(padded[:col_count])

        if not normalized_rows:
            for ev in events[:3]:
                base = [ev.get("date", ""), ev.get("headline", "")]
                base.extend([""] * max(0, col_count - len(base)))
                normalized_rows.append(base[:col_count])
        while len(normalized_rows) < 2:
            normalized_rows.append(list(normalized_rows[-1]))

        key_points = [p.strip() for p in (analysis.get("key_points") or []) if p and p.strip()][:5]
        if not key_points:
            key_points = [f"{ev.get('date', '')}: {ev.get('headline', '')}".strip(": ") for ev in events[:3]]
        key_points = [kp for kp in key_points if kp]
        if not key_points:
            key_points = ["待补充"]
        while len(key_points) < 3:
            key_points.append(key_points[-1])

        analysis_top = spine_y + connector_length + box_height + Inches(0.18)
        max_top = prs.slide_height - bottom_margin - Inches(1.6)
        guard_top = top_margin + Inches(0.3)
        if max_top < guard_top:
            max_top = guard_top
        if analysis_top > max_top:
            analysis_top = max(guard_top, max_top)
        analysis_height = max(Inches(1.6), prs.slide_height - bottom_margin - analysis_top)

        analysis_width = drawing_width
        left_panel_width = max(int(analysis_width * 0.58), int(min_left_panel))
        right_panel_width = analysis_width - left_panel_width - panel_gap
        if right_panel_width < min_right_panel:
            right_panel_width = min_right_panel
            left_panel_width = analysis_width - right_panel_width - panel_gap
        left_panel_left = left
        right_panel_left = left_panel_left + left_panel_width + panel_gap

        row_count = 1 + len(normalized_rows)
        table_shape = slide.shapes.add_table(
            row_count,
            col_count,
            left_panel_left,
            analysis_top,
            left_panel_width,
            analysis_height,
        )
        table = table_shape.table
        table.table_style = None

        for col_idx, header in enumerate(table_headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.text_frame.word_wrap = True
            header_paragraph = cell.text_frame.paragraphs[0]
            header_paragraph.font.size = Pt(14)
            header_paragraph.font.bold = True
            header_paragraph.font.name = theme.fonts["body"]
            header_paragraph.font.color.rgb = theme.colors["text"]
            header_paragraph.alignment = PP_ALIGN.CENTER
            cell.fill.background()
            _apply_cell_border(cell)

        for row_idx, row in enumerate(normalized_rows, start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = value
                cell.text_frame.word_wrap = True
                body_paragraph = cell.text_frame.paragraphs[0]
                body_paragraph.font.size = Pt(14)
                body_paragraph.font.name = theme.fonts["body"]
                body_paragraph.font.color.rgb = theme.colors["text"]
                body_paragraph.alignment = PP_ALIGN.LEFT
                cell.fill.background()
                _apply_cell_border(cell)

        separator_x = right_panel_left - panel_gap / 2
        separator = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            _to_emu(separator_x),
            _to_emu(analysis_top),
            _to_emu(separator_x),
            _to_emu(analysis_top + analysis_height),
        )
        separator.line.color.rgb = theme.colors["text"]
        separator.line.width = Pt(1)

        insights_box = slide.shapes.add_textbox(
            right_panel_left,
            analysis_top,
            right_panel_width,
            analysis_height,
        )
        tf_insights = insights_box.text_frame
        tf_insights.clear()
        tf_insights.word_wrap = True
        tf_insights.vertical_anchor = MSO_ANCHOR.TOP

        for idx, point in enumerate(key_points):
            paragraph = tf_insights.paragraphs[0] if idx == 0 else tf_insights.add_paragraph()
            paragraph.text = f"■ {point}"
            paragraph.font.size = Pt(14)
            paragraph.font.name = theme.fonts["body"]
            paragraph.font.color.rgb = theme.colors["text"]
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(6)
