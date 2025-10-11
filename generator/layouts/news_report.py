import re
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MsoShape
from ..ppt_builder import load_theme
from .utils import ensure_bg

FONT_NAME = "Microsoft YaHei"

# ---------- helper ----------
def val(obj, key, default=None):
    if obj is None:
        return default
    if isinstance(obj, dict):
        return obj.get(key, default)
    return getattr(obj, key, default)

def split_two_lines(text: str):
    t = (text or "").strip()
    if len(t) == 4:
        return t[:2], t[2:]
    return (t[:2] or "启示"), (t[2:] or "洞察")

def split_summary_explanation(text: str):
    t = (text or "").strip()
    if not t:
        return "要点", ""
    for sep in ['：', ':']:
        if sep in t:
            head, tail = t.split(sep, 1)
            return head.strip() or "要点", tail.strip()
    for sep in ['，', ',', '；', ';']:
        if sep in t:
            head, tail = t.split(sep, 1)
            return head.strip() or "要点", tail.strip()
    if len(t) > 8:
        return t[:6].strip(), t[6:].strip()
    return t, ""

def norm_bullets(x, mi=0, ma=12):
    if not x:
        return []
    if isinstance(x, list):
        arr = [str(s).strip() for s in x if str(s).strip()]
    else:
        parts = re.split(r"[；;。\.\n]+", str(x))
        arr = [p.strip(" 、，, ") for p in parts if p and p.strip()]
    return arr[:ma]

def no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass

# ---------- renderer ----------
def render(prs, routed_content):
    theme = load_theme()

    for slide_obj in routed_content.slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        ensure_bg(slide, theme)

        page_width = prs.slide_width
        margin_left = Inches(0.6)
        margin_right = Inches(0.6)

        # ===== 顶部：大标题 =====
        header_title = val(slide_obj, "header_title") or val(slide_obj, "title", "")
        title_box = slide.shapes.add_textbox(margin_left, Inches(0.22), page_width - margin_left - margin_right, Inches(1.0))
        tf_title = title_box.text_frame
        p = tf_title.paragraphs[0]
        p.text = header_title.strip()
        p.font.name = FONT_NAME
        p.font.size = Pt(theme.sizes["title_pt"])
        p.font.bold = True
        p.font.color.rgb = theme.colors["primary"]
        p.alignment = PP_ALIGN.LEFT

        # ===== 右上角品牌标签 =====
        brand_tag = str(val(slide_obj, "brand_tag", "CARI AI4News") or "CARI AI4News")
        tag_width = Inches(2.4)
        tag = slide.shapes.add_shape(
            MsoShape.RECTANGLE,
            page_width - margin_right - tag_width,
            Inches(0.22),
            tag_width,
            Inches(0.42),
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tag.line.fill.background()
        no_shadow(tag)
        tf = tag.text_frame
        brand_p = tf.paragraphs[0]
        brand_p.text = brand_tag
        brand_p.font.size = Pt(12)
        brand_p.font.name = FONT_NAME
        brand_p.font.color.rgb = RGBColor(0, 0, 0)
        brand_p.alignment = PP_ALIGN.CENTER

        # ===== 中部：Summary 灰框 =====
        summary = val(slide_obj, "summary", {}) or {}
        sum_label = val(summary, "label", "启示洞察")
        sum_bullets = norm_bullets(val(summary, "bullets"), mi=3, ma=6)

        block_h = Inches(0.9 + max(0, len(sum_bullets)-1) * 0.25)
        block_top = Inches(1.05)
        full_left = margin_left
        full_w = page_width - margin_left - margin_right

        bg = slide.shapes.add_shape(MsoShape.RECTANGLE, full_left, block_top, full_w, block_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(230, 230, 230)
        bg.line.fill.background()
        no_shadow(bg)

        red_w = Inches(1.2)
        red = slide.shapes.add_shape(MsoShape.RECTANGLE, full_left, block_top, red_w, block_h)
        red.fill.gradient()
        stops = red.fill.gradient_stops
        stops[0].color.rgb = RGBColor(180, 40, 40)
        stops[1].color.rgb = RGBColor(120, 0, 0)
        red.line.fill.background()
        no_shadow(red)

        l1, l2 = split_two_lines(sum_label)
        tf_red = red.text_frame
        tf_red.clear()
        tf_red.vertical_anchor = MSO_ANCHOR.MIDDLE
        fs = 16
        p1 = tf_red.paragraphs[0]
        p1.text = l1
        p1.font.bold = True
        p1.font.size = Pt(fs)
        p1.font.name = FONT_NAME
        p1.font.color.rgb = RGBColor(255, 255, 255)
        p1.alignment = PP_ALIGN.CENTER
        p2 = tf_red.add_paragraph()
        p2.text = l2
        p2.font.bold = True
        p2.font.size = Pt(fs)
        p2.font.name = FONT_NAME
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER

        sum_box = slide.shapes.add_textbox(full_left + red_w + Inches(0.12), block_top, full_w - red_w - Inches(0.12), block_h)
        tf = sum_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        for i, raw in enumerate(sum_bullets):
            summary_text, explanation = split_summary_explanation(raw)
            paragraph = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            paragraph.text = ""
            paragraph.font.size = Pt(12)
            paragraph.font.name = FONT_NAME
            paragraph.font.color.rgb = theme.colors["text"]

            run_bold = paragraph.add_run()
            run_bold.text = f"- {summary_text}："
            run_bold.font.bold = True
            run_bold.font.size = Pt(12)
            run_bold.font.name = FONT_NAME
            run_bold.font.color.rgb = theme.colors["text"]

            if explanation:
                run_rest = paragraph.add_run()
                run_rest.text = explanation.strip()
                run_rest.font.size = Pt(12)
                run_rest.font.name = FONT_NAME
                run_rest.font.color.rgb = theme.colors["text"]

            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(6)

        # ===== 下部：左右两列 =====
        grid_top = block_top + block_h + Inches(0.25)
        grid_gap = Inches(0.35)
        col_w = (full_w - grid_gap) / 2
        bottom_margin = theme.margins.get("bottom", Inches(0.6))
        extra_bottom_padding = Inches(0.3)
        grid_bottom = prs.slide_height - bottom_margin - extra_bottom_padding
        available_height = max(0, grid_bottom - grid_top)
        preferred_height = Inches(4.1)
        minimum_height = Inches(3.4)

        if available_height >= preferred_height:
            col_h = preferred_height
        elif available_height >= minimum_height:
            col_h = available_height
        else:
            col_h = minimum_height
            grid_bottom = grid_top + col_h

        left_col = val(slide_obj, "left", {}) or {}
        right_col = val(slide_obj, "right", {}) or {}

        def draw_column(col_left, col_data, fallback="列"):
            title_text = val(col_data, "title", fallback)
            sections = val(col_data, "sections", []) or []

            box = slide.shapes.add_shape(MsoShape.RECTANGLE, col_left, grid_top, col_w, col_h)
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(255, 255, 255)
            box.line.color.rgb = RGBColor(0, 102, 204)
            no_shadow(box)

            tag_w = Inches(3.5)
            tag_h = Inches(0.36)
            tag_x = col_left + (col_w - tag_w) / 2
            tag_y = grid_top - tag_h / 2
            tag = slide.shapes.add_shape(MsoShape.RECTANGLE, tag_x, tag_y, tag_w, tag_h)
            tag.fill.solid()
            tag.fill.fore_color.rgb = RGBColor(255, 255, 255)
            tag.line.fill.background()
            no_shadow(tag)
            tf_tag = tag.text_frame
            pt = tf_tag.paragraphs[0]
            pt.text = title_text
            pt.font.size = Pt(12)
            pt.font.bold = True
            pt.font.name = FONT_NAME
            pt.font.color.rgb = RGBColor(0, 0, 0)
            pt.alignment = PP_ALIGN.CENTER

            tfb = box.text_frame
            tfb.clear()
            tfb.word_wrap = True
            for idx, sec in enumerate(sections):
                subtitle = val(sec, "subtitle_bold", "") or f"小节{idx + 1}"
                bullets = norm_bullets(val(sec, "bullets"), mi=3, ma=6)

                psec = tfb.add_paragraph()
                psec.text = f"□ 【{subtitle}】"
                psec.font.size = Pt(12)
                psec.font.bold = True
                psec.font.name = FONT_NAME
                psec.font.color.rgb = RGBColor(0, 0, 0)
                psec.alignment = PP_ALIGN.LEFT
                psec.space_after = Pt(8)

                for b in bullets:
                    pb = tfb.add_paragraph()
                    pb.text = f"• {b}"
                    pb.font.size = Pt(12)
                    pb.font.name = FONT_NAME
                    pb.font.color.rgb = theme.colors["text"]
                    pb.alignment = PP_ALIGN.LEFT
                    pb.space_after = Pt(6)
            tfb.vertical_anchor = MSO_ANCHOR.TOP

        draw_column(full_left, left_col, "左列")
        draw_column(full_left + col_w + grid_gap, right_col, "右列")
