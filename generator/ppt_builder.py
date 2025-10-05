from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import yaml
from pathlib import Path
from typing import Dict, Any

HERE = Path(__file__).parent.parent

def _rgb(hex_str: str):
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2],16), int(hex_str[2:4],16), int(hex_str[4:6],16))

class Theme:
    def __init__(self, cfg: Dict[str, Any]):
        self.fonts = cfg["fonts"]
        self.sizes = cfg["sizes"]
        self.colors = {k: _rgb(v) for k, v in cfg["colors"].items()}
        m = cfg["margins"]
        self.margins = {k: Inches(v) for k, v in m.items()}
        self.logo = cfg.get("logo", {})
        self.bg = cfg.get("background", {})

def load_theme() -> Theme:
    with open(HERE / "templates" / "base_theme.yaml", "r", encoding="utf-8") as f:
        return Theme(yaml.safe_load(f))

def new_presentation(theme: Theme) -> Presentation:
    prs = Presentation()
    # 背景色/图在 python-pptx 需要逐页设（可在各 layout 中处理）
    return prs

def add_title(prs, theme: Theme, text: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = theme.margins["left"]; top = theme.margins["top"]
    width = Inches(10) - theme.margins["left"] - theme.margins["right"]
    height = Inches(1.0)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = theme.fonts["heading"]
    p.font.size = Pt(theme.sizes["title_pt"])
    p.font.bold = True
    p.font.color.rgb = theme.colors["primary"]   # 深红色
    p.alignment = PP_ALIGN.LEFT
    return slide